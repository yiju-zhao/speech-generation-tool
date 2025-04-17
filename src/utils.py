"""
Utility functions for the Digital Presenter project.
"""

import json
import toml
from pathlib import Path
from pptx import Presentation


# Load environment variables
def load_config():
    """Load configuration from config.toml file."""
    config_path = Path(__file__).parent.parent / "config" / "config.toml"

    if not config_path.exists():
        raise FileNotFoundError(
            f"Configuration file not found at {config_path}. "
            "Please create a config.toml file in the config directory."
        )

    try:
        config = toml.load(config_path)
        required_keys = ["knowledge_base_dir"]

        # Validate that all required keys are present
        missing_keys = [key for key in required_keys if key not in config]
        if missing_keys:
            raise KeyError(
                f"Missing required configuration keys: {', '.join(missing_keys)}. "
                "Please check your config.toml file."
            )

        return config
    except toml.TomlDecodeError as e:
        raise ValueError(f"Error parsing config.toml: {str(e)}")


def ensure_directory(directory):
    """Ensure a directory exists, create it if it doesn't."""
    Path(directory).mkdir(parents=True, exist_ok=True)
    return directory


def get_project_paths():
    """Get standard project paths."""
    base_dir = Path(__file__).parent.parent
    return {
        "base_dir": base_dir,
        "raw_dir": base_dir / "data" / "raw",
        "noted_dir": base_dir / "data" / "noted",
        "audio_dir": base_dir / "data" / "audio",
        "config_dir": base_dir / "config",
        "notebooks_dir": base_dir / "notebooks",
    }


def save_json(data, filename, directory=None):
    """Save data to a JSON file."""
    if directory:
        ensure_directory(directory)
        filepath = Path(directory) / filename
    else:
        filepath = Path(filename)

    with open(filepath, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2, ensure_ascii=False)

    return filepath


def load_json(filepath):
    """Load data from a JSON file."""
    with open(filepath, "r", encoding="utf-8") as f:
        return json.load(f)


def extract_transcripts_from_pptx(pptx_path):
    """
    Extract transcripts from the notes section of each slide in a PowerPoint file.

    Args:
        pptx_path (Path or str): Path to the PowerPoint file

    Returns:
        list: A list of dictionaries containing slide number and transcript
    """
    pptx_path = Path(pptx_path)

    # Load the presentation
    presentation = Presentation(pptx_path)

    # Extract transcripts from notes
    transcripts = []
    for i, slide in enumerate(presentation.slides, 1):
        # Get the content from the slide
        slide_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_content.append(shape.text.strip())

        original_content = "\n".join(slide_content) if slide_content else ""

        # Get transcript from notes
        transcript = ""
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text
            if notes_text.strip():
                transcript = notes_text.strip()

        slide_data = {
            "slide_number": i,
            "original_content": original_content,
            "transcript": transcript,
        }

        transcripts.append(slide_data)

    return transcripts


def write_transcripts_to_pptx(pptx_path, transcripts, output_base_dir=None):
    """Write transcripts to the notes of each slide in a PowerPoint presentation.

    Args:
        pptx_path: Path to the input PowerPoint file
        transcripts: List of dictionaries with transcript information
        output_base_dir: Directory to save the output file (optional)

    Returns:
        Path to the saved PowerPoint file with transcripts
    """
    # Load the presentation
    presentation = Presentation(pptx_path)

    # Create a dictionary mapping slide numbers to transcripts for easier lookup
    slide_transcript_map = {data["slide_number"]: data for data in transcripts}
    
    # Write transcripts to the notes of each slide
    for i, slide in enumerate(presentation.slides, 1):  # Start from 1 to match slide_number
        # Check if we have a transcript for this slide
        if i in slide_transcript_map:
            # Get the transcript for this slide
            slide_data = slide_transcript_map[i]

            # Create a notes slide if it doesn't exist
            if not slide.has_notes_slide:
                slide.notes_slide

            # Set the notes text to the transcript
            slide.notes_slide.notes_text_frame.text = slide_data["transcript"]

    # Determine output path
    if output_base_dir is None:
        output_path = Path(str(pptx_path).replace(".pptx", "_noted.pptx"))
    else:
        output_base_dir = Path(output_base_dir)
        output_base_dir.mkdir(parents=True, exist_ok=True)
        output_path = output_base_dir / f"{Path(pptx_path).stem}_noted.pptx"

    # Save the presentation
    presentation.save(output_path)

    return str(output_path)


def extract_slide_content(slide):
    """Extract all text content from a slide."""
    slide_content = []

    # Process all shapes in the slide
    for shape in slide.shapes:
        # Handle direct text attributes
        if hasattr(shape, "text") and shape.text.strip():
            slide_content.append(shape.text.strip())

        # Handle tables
        if shape.has_table:
            table = shape.table
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    slide_content.append(" | ".join(row_text))

        # Handle text frames (which might contain paragraphs with runs)
        if hasattr(shape, "text_frame"):
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                paragraph_text = paragraph.text.strip()
                if paragraph_text:
                    slide_content.append(paragraph_text)

    return "\n".join(slide_content) if slide_content else None
