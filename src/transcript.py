"""
Transcript generation module for Digital Presenter.
"""

from pptx import Presentation

from .utils import (
    load_config,
    ensure_directory,
    get_project_paths,
    write_transcripts_to_pptx,
    extract_slide_content
)
from .llm import get_llm_provider


def generate_transcript(slide_content, previous_transcripts=None, llm_client=None, target_language=None, model=None):
    """Generate a transcript for a slide using LLM."""
    # Add all previous slides' transcripts as context if available
    context = ""
    if previous_transcripts:
        context = f"{previous_transcripts}\n"
    
    # Create a unified prompt that handles both cases
    prompt = f"""
    **Role**: Technical Speech Architect  
    **Objective**: Create TTS-ready transcripts with uncompromised accuracy  

    **Core Directives**:  
    1. Preserve _exact_ terminology + numerical precision  
    2. Optimize complex concepts for vocal clarity  
    3. Maintain academic rigor in spoken format  

    **Input Sources**:  
    - Content: {slide_content}  
    - Context: {context}  

    **Output Specifications**:  
    `Language`: {target_language.upper() if target_language else "ENGLISH"}  
    `Format`: Pure speech text (no markup/annotations)  

    **Technical Protocols**:  
    ✓ Term fidelity: Use source terms verbatim  
    ✓ Data integrity: Exact values + full SI units  
    ✓ First-use acronyms: "Advanced Reactor Concept (ARC)"  
    ✓ Concept links: "This correlates directly with <previous_term>"  
    ✓ Transition logic: "Consequently...", "The evidence establishes..."  

    **Prohibited Elements**:  
    - Approximations ("~", "about")  
    - Subjective commentary ("surprisingly")  
    - Non-verbal cues (parentheticals, pauses)  
    - Casual paraphrasing  

    **Structural Guidelines**:  
    1. **Opening**: Bridge from {context} using exact technical references  
    2. **Flow**:  
    - State quantitative relationships verbally  
    - Use 3-tier vocal hierarchy: Concept → Mechanism → Evidence  
    - Employ precision transitions: "The critical progression..."  
    3. **Complex Data**:  
    - "Three components emerge: (1) <term>, (2) <term>..."  
    - "Sequential analysis reveals: First... Second... Crucially..."  

    **Safety Imperatives**:  
    ‼ Preserve safety-critical wording verbatim  
    ‼ Verbalize technical caveats explicitly  
    ‼ If term conflict: Prioritize slide terminology  
    
    **FORMAT RESTRICTION (CRITICAL)**:
    1. Output ONLY the actual transcript text
    2. DO NOT include any introductory phrases like "Here is the transcript:" or "以下是转录内容:"
    3. DO NOT include any explanations or notes after the transcript
    4. DO NOT use quotation marks to wrap the entire transcript
    5. Return ONLY the pure speech text that would be read aloud
    """
    
    return llm_client.generate(prompt, model)


def process_presentation(pptx_path, output_base_dir=None, target_language=None, model=None, llm_provider=None, slides=None):
    """Process a presentation and generate transcripts for specified slides."""
    # Load the presentation
    presentation = Presentation(pptx_path)
    
    # Set up output directory structure
    if not output_base_dir:
        paths = get_project_paths()
        output_base_dir = paths["noted_dir"]
    ensure_directory(output_base_dir)

    # Initialize the appropriate LLM client based on model choice
    config = load_config()
    llm_client = get_llm_provider(model, config, provider=llm_provider)

    # Process each slide
    transcripts = []
    previous_transcripts = []
    for idx, slide in enumerate(presentation.slides, 1):
        if slides and idx not in slides:
            print(f"Skipping slide {idx} (not in specified slides)")
            continue
        
        # Extract all text from the slide
        slide_content = extract_slide_content(slide)
        
        # Check for unpolished notes in the slide
        unpolished_notes = ""
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text
            if notes_text.strip():
                unpolished_notes = notes_text.strip()
                print(f"Found unpolished notes for slide {idx}")

        if not slide_content and not unpolished_notes:
            print(f"No content found for slide {idx}, skipping.")
            continue

        # Combine slide content and unpolished notes if available
        combined_content = []
        if slide_content:
            combined_content.append(f"SLIDE CONTENT:\n{slide_content}")
        if unpolished_notes:
            combined_content.append(f"UNPOLISHED NOTES (USE AS GUIDE):\n{unpolished_notes}")
        
        # Generate transcript with combined content
        transcript = generate_transcript(
            "\n\n".join(combined_content),
            previous_transcripts,
            llm_client,
            target_language,
            model
        )

        # Store results
        slide_data = {
            "slide_number": idx,
            "original_content": slide_content,
            "unpolished_notes": unpolished_notes,
            "transcript": transcript,
        }
        transcripts.append(slide_data)
        previous_transcripts = previous_transcripts + [transcript]

        print(f"Processed slide {idx}")

    # Write transcripts back to the PPTX file
    noted_pptx_path = write_transcripts_to_pptx(pptx_path, transcripts, output_base_dir)
    print(f"PPTX with transcripts has been saved to {noted_pptx_path}")

    return noted_pptx_path


def main():
    """Main function to run the transcript generator."""
    from config.parse_args import parse_args
    
    # Get arguments from parse_args and filter what we need
    args = parse_args()
    use_storm = args.use_storm
    language = args.language
    model = args.model
    llm_provider = args.llm_provider
    
    paths = get_project_paths()
    ensure_directory(paths["input_dir"])
    ensure_directory(paths["output_dir"])

    # Find PowerPoint files in the input directory
    pptx_files = list(paths["input_dir"].glob("*.pptx"))

    if not pptx_files:
        print(f"No PowerPoint files found in {paths['input_dir']}")
        return

    # Import storm module only if needed
    if use_storm:
        try:
            from .transcript_storm import process_presentation_with_storm
            print("Using Storm-enhanced approach for transcript generation...")
        except ImportError:
            print("Warning: Storm module not found. Falling back to standard approach.")
            use_storm = False

    # Process each PowerPoint file
    for pptx_file in pptx_files:
        print(f"Processing {pptx_file.name}...")
        
        output_dir = paths["output_dir"] / pptx_file.stem
        
        if use_storm:
            print("Using enhanced verification workflow to reduce hallucinations.")
            process_presentation_with_storm(pptx_file, output_dir, language, model, llm_provider=llm_provider)
        else:
            print("Using standard transcript generation workflow.")
            process_presentation(pptx_file, output_dir, language, model, llm_provider=llm_provider)
        
        print(f"Transcript generation complete for {pptx_file.name}")
        print("Check the output directory for the PPTX file with transcripts in the notes section.")


if __name__ == "__main__":
    main()
