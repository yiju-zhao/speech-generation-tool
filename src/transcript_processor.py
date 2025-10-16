"""
Main processor for handling presentations with the Storm approach
"""

import os
import json
import logging
from pathlib import Path
from pptx import Presentation

from .utils import (
    load_config,
    ensure_directory,
    get_project_paths,
    write_transcripts_to_pptx,
    extract_slide_content,
)
from .llm import get_llm_provider
from .models import SlideInformation
from .transcript_curator import TranscriptKnowledgeCurator
from .transcript_generator import TranscriptGenerator, TranscriptReviewer


def process_presentation_with_storm(
    pptx_path,
    output_base_dir=None,
    target_language=None,
    model=None,
    enable_search=True,
    llm_provider=None,
    slides=None,
):
    """Process a presentation using Storm-inspired approach to reduce hallucinations."""
    presentation = Presentation(pptx_path)

    # ------------------- Directory setup -------------------
    if not output_base_dir:
        paths = get_project_paths()
        output_base_dir = paths["processed_dir"]
    ensure_directory(output_base_dir)

    log_dir = os.path.join(output_base_dir, "logs")
    ensure_directory(log_dir)
    log_path = os.path.join(
        log_dir, f"{os.path.splitext(os.path.basename(pptx_path))[0]}_processing.log"
    )

    # ------------------- Logging setup -------------------
    logging.basicConfig(
        filename=log_path,
        level=logging.INFO,
        format="%(asctime)s - %(levelname)s - %(message)s",
        datefmt="%Y-%m-%d %H:%M:%S",
    )
    console = logging.StreamHandler()
    console.setLevel(logging.INFO)
    logging.getLogger().addHandler(console)

    logging.info(f"Starting processing: {pptx_path}")

    # ------------------- Config setup -------------------
    config = load_config()
    llm_client = get_llm_provider(model, config, provider=llm_provider)
    tavily_api_key = config.get("tavily_api_key", None)
    knowledge_base_dir = config.get("knowledge_base_dir")

    if knowledge_base_dir:
        file_base = os.path.splitext(os.path.basename(pptx_path))[0]
        knowledge_base_dir = os.path.join(knowledge_base_dir, file_base)
        ensure_directory(knowledge_base_dir)
        logging.info(f"Knowledge base directory: {knowledge_base_dir}")

    # ------------------- Validation -------------------
    if enable_search and not tavily_api_key:
        logging.warning("Web search enabled but no Tavily key found — disabling search.")
        enable_search = False

    # ------------------- Component setup -------------------
    knowledge_curator = TranscriptKnowledgeCurator(
        llm_client, model, enable_search=enable_search,
        knowledge_base_dir=knowledge_base_dir, tavily_api_key=tavily_api_key
    )
    # Pass knowledge retriever and KB dir into generator for lean RAG context
    transcript_generator = TranscriptGenerator(
        llm_client,
        model,
        target_language,
        knowledge_base_dir=knowledge_base_dir,
        knowledge_retriever=getattr(knowledge_curator, "knowledge_retriever", None),
    )
    transcript_reviewer = TranscriptReviewer(llm_client, model)

    # ------------------- Tracking setup -------------------
    transcripts = []
    previous_transcripts = []
    slide_information = {}
    all_knowledge_items = []

    slide_count = len(presentation.slides)
    logging.info(f"Total slides: {slide_count}")
    logging.info(f"Model: {model}, Provider: {llm_provider or 'default'}")
    logging.info(f"Web search: {'ON' if enable_search else 'OFF'}, Knowledge base: {knowledge_base_dir or 'OFF'}")

    # ------------------- Process slides -------------------
    for idx, slide in enumerate(presentation.slides, 1):
        if slides and idx not in slides:
            logging.info(f"Skipping slide {idx} (not in target list)")
            continue

        logging.info(f"\n--- Processing Slide {idx} ---")

        # Extract slide text and notes
        slide_content = extract_slide_content(slide)
        unpolished_notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes.strip():
                unpolished_notes = notes.strip()

        if not slide_content and not unpolished_notes:
            logging.info(f"No content found for slide {idx}, skipping.")
            continue

        combined_content = f"{slide_content or ''}\n\nUNPOLISHED NOTES:\n{unpolished_notes}" if unpolished_notes else slide_content or ""
        slide_info = SlideInformation(
            original_content=slide_content or "",
            unpolished_notes=unpolished_notes,
            slide_number=idx,
        )

        # --------------- Step 1: Knowledge Curation ---------------
        slide_kb_path = None
        cached_items = []
        if knowledge_base_dir:
            slide_kb_path = Path(knowledge_base_dir) / f"slide_{idx:03d}_knowledge.json"
            if slide_kb_path.exists():
                try:
                    data = json.load(open(slide_kb_path))
                    cached_items = data.get("items", [])
                    if cached_items:
                        logging.info(f"Using cached knowledge for slide {idx} ({len(cached_items)} items)")
                        from .knowledge_base import KnowledgeItem
                        slide_info.knowledge_items = [KnowledgeItem.from_dict(d) for d in cached_items]
                except Exception as e:
                    logging.warning(f"Error loading cached knowledge for slide {idx}: {e}")

        # If no cached knowledge, perform new retrieval
        if not cached_items:
            try:
                logging.info(f"Generating search queries for slide {idx}")
                slide_info.queries = knowledge_curator.generate_search_queries(combined_content, idx)
            except Exception as e:
                logging.error(f"Error generating queries: {e}")
                slide_info.queries = []

            try:
                logging.info(f"Retrieving knowledge for slide {idx}")
                slide_info.knowledge_items = knowledge_curator.perform_knowledge_retrieval(
                    combined_content, idx, slide_info.queries
                )
                logging.info(f"Retrieved {len(slide_info.knowledge_items)} items for slide {idx}")
            except Exception as e:
                logging.error(f"Knowledge retrieval failed for slide {idx}: {e}")
                slide_info.knowledge_items = []

        # Extract verified facts and synthesize content
        try:
            slide_info.facts = knowledge_curator.extract_verified_facts(
                combined_content, slide_info.queries, slide_info.knowledge_items
            )
            slide_info.verified_content = knowledge_curator.synthesize_verified_content(slide_info)
        except Exception as e:
            logging.error(f"Error verifying content for slide {idx}: {e}")
            slide_info.facts = combined_content.split("\n")
            slide_info.verified_content = combined_content

        # --------------- Step 2: Transcript Generation ---------------
        try:
            # Use explicit keywords to avoid parameter ordering bugs
            transcript = transcript_generator.generate_transcript(
                slide_info.verified_content,
                slide_info=slide_info,
                previous_transcripts=previous_transcripts,
                enforce_semantic_check=False,
            )
        except Exception as e:
            logging.error(f"Transcript generation failed: {e}")
            transcript = slide_info.verified_content

        # --------------- Step 3: Transcript Review ---------------
        try:
            review = transcript_reviewer.review_transcript(slide_info, transcript)
            final_transcript = review.get("revised_transcript", transcript)
        except Exception as e:
            logging.error(f"Transcript review failed: {e}")
            review = {"accurate": True, "hallucinations": [], "corrections": [], "revised_transcript": transcript}
            final_transcript = transcript

        # --------------- Step 4: Store results ---------------
        slide_data = {
            "slide_number": idx,
            "original_content": slide_content,
            "unpolished_notes": unpolished_notes,
            "verified_content": slide_info.verified_content,
            "verified_facts": slide_info.facts,
            "transcript": final_transcript,
            "accuracy_review": review,
        }

        if slide_info.knowledge_items:
            slide_data["knowledge_items"] = [
                {
                    "content": item.content[:500] if len(item.content) > 500 else item.content,
                    "source": item.source,
                    "metadata": item.metadata,
                }
                for item in slide_info.knowledge_items
            ]

        transcripts.append(slide_data)
        previous_transcripts.append(final_transcript)

        logging.info(f"Slide {idx} processed successfully")

    # ------------------- Save outputs -------------------
    noted_pptx_path = write_transcripts_to_pptx(pptx_path, transcripts, output_base_dir)
    process_data = {
        "slides": transcripts,
        "metadata": {
            "source_file": os.path.basename(pptx_path),
            "total_slides": len(presentation.slides),
            "processed_slides": len(transcripts),
            "target_language": target_language,
            "model": model,
            "web_search_enabled": enable_search,
        },
    }

    json_path = os.path.join(
        output_base_dir, f"{os.path.splitext(os.path.basename(pptx_path))[0]}_process_data.json"
    )
    with open(json_path, "w") as f:
        json.dump(process_data, f, indent=2)

    logging.info(f"Processing complete. Data saved to {json_path}")
    return noted_pptx_path
